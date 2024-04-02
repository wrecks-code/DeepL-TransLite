import os
import webbrowser

import pptx

from deepl_pptx_translator import apiHandler, configHandler, guiHandler, textHandler

include_subdirectories = True
total_files = 0
processed_files = 0
translated_files_count = 0


def open_output_folder_in_explorer():
    # Convert the relative path to an absolute path
    absolute_path = os.path.abspath(
        os.path.join(os.getcwd(), configHandler.output_path)
    )
    # Open the "Output" folder using the default file explorer
    webbrowser.open(absolute_path)


def process_pptx_files(input_folder_local):
    global total_files, processed_files
    total_files = sum(
        1 for file in os.listdir(input_folder_local) if file.endswith(".pptx")
    )
    processed_files = 0

    # Check if the output folder exists
    if not os.path.exists(configHandler.output_path):
        print(f"The folder '{configHandler.output_path}' does not exist. Creating...")
        os.makedirs(configHandler.output_path)

    # Initialize a stack with the root directory
    stack = [input_folder_local]

    # Process directories until the stack is empty
    while stack:
        current_dir = stack.pop()

        # Iterate through files and directories in the current directory
        for item in os.listdir(current_dir):
            item_path = os.path.join(current_dir, item)

            if os.path.isdir(item_path) and include_subdirectories:
                # If item is a directory and include_subdirectories is True, add it to the stack for processing
                stack.append(item_path)
            elif item.endswith(".pptx"):
                # If item is a .pptx file, process it
                translate_presentation(item_path)
                processed_files += 1
                progress = processed_files / total_files * 100
                guiHandler.progress_label.config(
                    text=str(processed_files) + " von " + str(total_files)
                )
                guiHandler.progress_bar["value"] = progress
                guiHandler.progress_bar.update_idletasks()


def count_characters_in_presentation(input_path_local):
    total_characters = 0
    presentation_local = pptx.Presentation(input_path_local)
    presentation_name = os.path.basename(input_path_local)

    for slide in presentation_local.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                total_characters += len(shape.text)

    print(f"Total characters in {presentation_name} is {total_characters}")
    return total_characters


def count_characters_in_folder(folder_path):
    total_characters = 0

    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            file_path = os.path.join(folder_path, filename)
            total_characters += count_characters_in_presentation(file_path)

    print(f"Total characters in folder {folder_path} is {total_characters}")
    return total_characters


def translate_presentation(input_path_translate):
    global translated_files_count  # Add a global variable to keep track of the count
    translated_files_count += 1  # Increment the count for each translated file

    presentation = pptx.Presentation(input_path_translate)
    total_slides = sum(1 for _ in presentation.slides)
    count_characters_in_presentation(input_path_translate)

    # Translate the entire title using Deepl
    translated_title = apiHandler.translate_text_w_deepl(
        os.path.basename(input_path_translate)
    )

    pptx_file_path = os.path.join(configHandler.output_path, f"{translated_title}")

    # Initialize a counter for the processed slides
    processed_slides = 0

    # Iterate through each slide in the presentation
    for slide_count, slide in enumerate(presentation.slides):
        print(f"Slide (Folie): {slide_count}")

        # Increment the processed slides counter
        processed_slides += 1

        guiHandler.progress_label.config(text="")

        # Calculate the progress and update the progress bar
        progress = processed_slides / total_slides * 100
        guiHandler.progress_bar["value"] = progress
        guiHandler.progress_bar.update_idletasks()

        # Define a function to process shapes recursively
        def process_shapes(shapes):
            for shape_count, shape in enumerate(shapes):
                print(f"Shape (Form): {shape_count}")

                # Check if the shape contains text
                if shape.has_text_frame and shape.text:
                    # Iterate through each paragraph and run in the shape
                    for para_count, paragraph in enumerate(shape.text_frame.paragraphs):
                        print(f"Paragraph (Satz): {para_count}")
                        sentence_to_translate = ""
                        translated_sentence = ""

                        # Iterate through each run (text segment) in the paragraph
                        for run in paragraph.runs:
                            print(f"Run (Satzabschnitt): '{run.text}'")
                            # Concatenate text segments to form a sentence
                            sentence_to_translate += textHandler.add_plus(run.text)

                        print("Sentence to translate: " + sentence_to_translate)
                        translated_sentence = apiHandler.translate_text_w_deepl(
                            sentence_to_translate
                        )
                        print("translated sentence: " + translated_sentence)
                        textHandler.assign_segments_to_runs(
                            paragraph,
                            textHandler.split_text_with_marker(translated_sentence),
                        )

                # Recursively process nested shapes
                if hasattr(shape, "shapes"):
                    process_shapes(shape.shapes)

        # Start processing from the top-level shapes
        process_shapes(slide.shapes)

    # Save the translated presentation
    presentation.save(pptx_file_path)
