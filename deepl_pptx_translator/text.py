import os
import pptx
import docx
from langdetect import detect, LangDetectException
from deepl_pptx_translator import config_handler, gui_handler


def add_plus(text) -> str:
    return config_handler.marker_char + text + config_handler.marker_char


def split_text_with_marker(text) -> list:
    # print("splitting " + text)
    segments = [
        segment for segment in text.split(config_handler.marker_char) if segment
    ]
    # print("splitted version: " + str(segments))
    return segments


def assign_segments_to_runs(paragraph, segments):
    # print("size of paragraph.runs: " + str(len(paragraph.runs)))
    # print("size of segments: " + str(len(segments)))

    # Extend the segments list with empty strings if it's shorter than paragraph.runs
    if len(paragraph.runs) > len(segments):
        segments.extend([""] * (len(paragraph.runs) - len(segments)))

    # Iterate through the runs in the paragraph and assign segments to each run
    for i, run in enumerate(paragraph.runs):
        run.text = segments[i]


def detect_language(selected_path) -> str:
    def detect_language_for_file(file_path) -> str:
        if file_path.endswith(".docx"):
            return detect_docx_language(file_path)
        elif file_path.endswith(".pptx"):
            return detect_pptx_language(file_path)
        return None

    language_code = None

    if os.path.isdir(selected_path):
        for root, dirs, files in os.walk(selected_path):
            for file in files:
                file_path = os.path.join(root, file)
                language_code = detect_language_for_file(file_path)
                if language_code:  # If a language code is found, no need to continue
                    break
            if (
                language_code
            ):  # If a language code is found, break out of the outer loop
                break
    else:
        language_code = detect_language_for_file(selected_path)

    if language_code == "DE":
        gui_handler.TARGET_CHOSEN_LANG.set("EN-GB")
    elif language_code == "EN-GB":
        gui_handler.TARGET_CHOSEN_LANG.set("DE")

    return language_code


def detect_docx_language(docx_path) -> str:
    try:
        document = docx.Document(docx_path)
        return_text = " ".join(
            paragraph.text for paragraph in document.paragraphs if paragraph.text
        )
        return get_language_code(return_text)
    except Exception as e:
        print(f"Error processing DOCX file: {e}")
        return None


def detect_pptx_language(pptx_path) -> str:
    try:
        presentation = pptx.Presentation(pptx_path)
        return_text = " ".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        return get_language_code(return_text)
    except Exception as e:
        print(f"Error processing PPTX file: {e}")
        return None


def get_language_code(_text) -> str:
    try:
        lang = detect(_text)
        if lang == "en":
            return "EN-GB"
        elif lang == "de":
            return "DE"
        else:
            return lang.upper()
    except LangDetectException:
        return None
