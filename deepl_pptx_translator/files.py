import os
import webbrowser

import pptx
import docx

from deepl_pptx_translator import api, config_handler, gui_handler, text

INCLUDE_SUBDIRS = True
TOTAL_FILES = 0
PROCESSED_FILES = 0
TRANSLATED_FILES_COUNT = 0

# pylint: disable=W0603
# pylint: disable=W0640


def open_output_folder_in_explorer():
    # Convert the relative path to an absolute path
    absolute_path = os.path.abspath(
        os.path.join(os.getcwd(), config_handler.output_path)
    )
    # Open the "Output" folder using the default file explorer
    webbrowser.open(absolute_path)


def process_presentation_and_document_files(input_folder_local):
    global TOTAL_FILES, PROCESSED_FILES
    # Count total .pptx and .docx files
    TOTAL_FILES = sum(
        1
        for file in os.listdir(input_folder_local)
        if file.endswith((".pptx", ".docx"))
    )
    PROCESSED_FILES = 0

    # Check if the output folder exists
    if not os.path.exists(config_handler.output_path):
        print(f"The folder '{config_handler.output_path}' does not exist. Creating...")
        os.makedirs(config_handler.output_path)

    # Initialize a stack with the root directory
    stack = [input_folder_local]

    # Process directories until the stack is empty
    while stack:
        current_dir = stack.pop()

        # Iterate through files and directories in the current directory
        for item in os.listdir(current_dir):
            item_path = os.path.join(current_dir, item)

            if os.path.isdir(item_path) and INCLUDE_SUBDIRS:
                # If item is a directory and include_subdirectories is True, add it to the stack for processing
                stack.append(item_path)
            elif item.endswith(".pptx"):
                # If item is a .pptx file, process it
                translate_presentation(item_path)
                PROCESSED_FILES += 1
            elif item.endswith(".docx"):
                # If item is a .docx file, process it
                translate_document(item_path)
                PROCESSED_FILES += 1

            # Update progress
            if item.endswith((".pptx", ".docx")):
                progress = PROCESSED_FILES / TOTAL_FILES * 100
                gui_handler.PROGRESS_LABEL.config(
                    text=f"{PROCESSED_FILES} von {TOTAL_FILES}"
                )
                gui_handler.PROGRESS_BAR["value"] = progress
                gui_handler.PROGRESS_BAR.update_idletasks()


def count_characters_in_file(input_path_local) -> int:
    total_characters = 0
    file_extension = os.path.splitext(input_path_local)[1].lower()
    file_name = os.path.basename(input_path_local)

    if file_extension == ".pptx":
        # Process the .pptx file
        presentation_local = pptx.Presentation(input_path_local)
        for slide in presentation_local.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    total_characters += len(shape.text)
    elif file_extension == ".docx":
        # Process the .docx file
        doc = docx.Document(input_path_local)
        for paragraph in doc.paragraphs:
            total_characters += len(paragraph.text)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

    print(f"Total characters in {file_name} is {total_characters}")
    return total_characters


def count_characters_in_folder(folder_path) -> int:
    total_characters = 0

    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx") or filename.endswith(".docx"):
            file_path = os.path.join(folder_path, filename)
            total_characters += count_characters_in_file(file_path)

    print(f"Total characters in folder {folder_path} is {total_characters}")
    return total_characters


def count_docx_pptx_files_in_folder(folder_path) -> int:

    if INCLUDE_SUBDIRS:
        count = 0
        for root, dirs, files in os.walk(folder_path):
            count += len([name for name in files if name.endswith((".docx", ".pptx"))])
        return count
    else:
        return len(
            [
                name
                for name in os.listdir(folder_path)
                if name.endswith((".docx", ".pptx"))
                and os.path.isfile(os.path.join(folder_path, name))
            ]
        )


def translate_presentation(input_path_translate):
    global TRANSLATED_FILES_COUNT  # Add a global variable to keep track of the count
    TRANSLATED_FILES_COUNT += 1  # Increment the count for each translated file

    presentation = pptx.Presentation(input_path_translate)
    total_slides = sum(1 for _ in presentation.slides)
    count_characters_in_file(input_path_translate)

    # Translate the entire title using Deepl
    translated_title = api.translate_text_w_deepl(
        os.path.basename(input_path_translate)
    )

    pptx_file_path = os.path.join(config_handler.output_path, f"{translated_title}")

    # Initialize a counter for the processed slides
    processed_slides = 0

    # Iterate through each slide in the presentation
    for slide_count, slide in enumerate(presentation.slides):
        print(f"Slide (Folie): {slide_count}")

        # Increment the processed slides counter
        processed_slides += 1

        gui_handler.PROGRESS_LABEL.config(text="")

        # Calculate the progress and update the progress bar
        progress = processed_slides / total_slides * 100
        gui_handler.PROGRESS_BAR["value"] = progress
        gui_handler.PROGRESS_BAR.update_idletasks()

        # Define a function to process shapes recursively
        def process_shapes(shapes):
            for shape_count, shape in enumerate(shapes):
                print(f"Shape (Form): {shape_count}")

                # Check if the shape contains text
                if shape.has_text_frame and shape.text:
                    # Iterate through each paragraph and run in the shape
                    for para_count, paragraph in enumerate(shape.text_frame.paragraphs):
                        print(f"Paragraph (Satz): {para_count}")
                        sentence_to_translate = ""
                        translated_sentence = ""

                        # Iterate through each run (text segment) in the paragraph
                        for run in paragraph.runs:
                            print(f"Run (Satzabschnitt): '{run.text}'")
                            # Concatenate text segments to form a sentence
                            sentence_to_translate += text.add_plus(run.text)

                        print("Sentence to translate: " + sentence_to_translate)
                        translated_sentence = api.translate_text_w_deepl(
                            sentence_to_translate
                        )
                        print("translated sentence: " + translated_sentence)
                        text.assign_segments_to_runs(
                            paragraph,
                            text.split_text_with_marker(translated_sentence),
                        )

                # Recursively process nested shapes
                if hasattr(shape, "shapes"):
                    process_shapes(shape.shapes)

        # Start processing from the top-level shapes
        process_shapes(slide.shapes)

    # Save the translated presentation
    presentation.save(pptx_file_path)


def translate_document(input_path_translate):
    global TRANSLATED_FILES_COUNT
    TRANSLATED_FILES_COUNT += 1  # Increment the count for each translated file

    document = docx.Document(input_path_translate)
    total_paragraphs = sum(1 for _ in document.paragraphs)

    # Translate the entire title using Deepl
    translated_title = api.translate_text_w_deepl(
        os.path.basename(input_path_translate)
    )

    docx_file_path = os.path.join(config_handler.output_path, f"{translated_title}")

    # Initialize a counter for the processed paragraphs
    processed_paragraphs = 0

    # Define a function to process runs within paragraphs
    def process_runs(paragraph):
        sentence_to_translate = ""
        translated_sentence = ""
        for run in paragraph.runs:
            print(f"Run (Satzabschnitt): '{run.text}'")
            # Concatenate text segments to form a sentence
            sentence_to_translate += text.add_plus(run.text)

        print("Sentence to translate: " + sentence_to_translate)
        translated_sentence = api.translate_text_w_deepl(sentence_to_translate)
        print("Translated sentence: " + translated_sentence)
        text.assign_segments_to_runs(
            paragraph, text.split_text_with_marker(translated_sentence)
        )

    # Iterate through each paragraph in the document
    for para_count, paragraph in enumerate(document.paragraphs):
        print(f"Paragraph (Satz): {para_count}")

        # Increment the processed paragraphs counter
        processed_paragraphs += 1

        gui_handler.PROGRESS_LABEL.config(text="")

        # Calculate the progress and update the progress bar
        progress = processed_paragraphs / total_paragraphs * 100
        gui_handler.PROGRESS_BAR["value"] = progress
        gui_handler.PROGRESS_BAR.update_idletasks()

        # Process runs within the paragraph
        process_runs(paragraph)

    # Save the translated document
    document.save(docx_file_path)


def check_file_type(file_path) -> str:
    _, file_extension = os.path.splitext(file_path)
    if file_extension.lower() == ".docx":
        return "docx"
    elif file_extension.lower() == ".pptx":
        return "pptx"
    else:
        return None
