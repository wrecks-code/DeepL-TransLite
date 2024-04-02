from pptx import Presentation


# Testing this
def extract_all_text(presentation_path):
    """
    Extracts all text content from a PowerPoint presentation, including nested shapes.

    Args:
      presentation_path: Path to the PowerPoint presentation file.

    Returns:
      A list containing all extracted text elements.
    """
    prs = Presentation(presentation_path)
    all_text = []

    def explore_shapes(shapes):
        for shape in shapes:
            # Check for text frame
            if shape.has_text_frame:
                # Extract text from paragraphs
                for paragraph in shape.text_frame.paragraphs:
                    all_text.append(paragraph.text)
            # Recursively explore groups and subgroups (nested shapes)
            if shape.type == Presentation.enum.shapes.MsoShapeType.GROUP:
                explore_shapes(shape.group)  # Group property holds nested shapes

    # Iterate through all slides and their shapes (including nested)
    for slide in prs.slides:
        explore_shapes(slide.shapes)

    return all_text


# Example usage
presentation_file = "test_scripts/testfile.pptx"
all_text_content = extract_all_text(presentation_file)

# Print all extracted text (optional)
for text in all_text_content:
    print(text)
